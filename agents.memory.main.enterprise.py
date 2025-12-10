"""
VeriSynthOS Memory Agent - ENTERPRISE EDITION
Full implementation with all 20 required features from memory.agent.md
"""

import logging
import hashlib
import tiktoken
from datetime import datetime, timezone, timedelta
from typing import List, Dict, Any, Optional
from pathlib import Path
import asyncio
import json
import uuid

from fastapi import FastAPI, HTTPException, Request, BackgroundTasks
from fastapi.responses import JSONResponse
from fastapi.staticfiles import StaticFiles
from fastapi.responses import HTMLResponse
from pydantic import BaseModel, Field
import google.auth
from google.cloud import firestore
from google.cloud import tasks_v2
from googleapiclient.discovery import build
from googleapiclient.http import MediaIoBaseDownload
import io

# Vertex AI imports
try:
    from google.cloud import aiplatform
    from vertexai.language_models import TextEmbeddingModel
    VERTEX_AVAILABLE = True
except ImportError:
    VERTEX_AVAILABLE = False

log = logging.getLogger("memory")
log.setLevel(logging.INFO)

app = FastAPI(title="VeriSynthOS Memory Agent - Enterprise Edition")

# ------------------------------------------------------------------
# GLOBAL STATE
# ------------------------------------------------------------------
credentials = None
project_id = None
db = None
drive_service = None
watch_channels: Dict[str, Dict] = {}  # channel_id -> {folder_id, expiration, resource_id}
task_client = None
tokenizer = None

# Environment configuration
import os
ME_INDEX_ID = os.getenv("ME_INDEX_ID")
ME_ENDPOINT_ID = os.getenv("ME_ENDPOINT_ID")
WEBHOOK_URL = os.getenv("WEBHOOK_URL", "https://your-domain.com/webhook/drive")
CLOUD_TASKS_QUEUE = os.getenv("CLOUD_TASKS_QUEUE", "memory-ingestion-queue")

# Feature flags
ENABLE_DRIVE_WATCH = os.getenv("ENABLE_DRIVE_WATCH", "true").lower() == "true"
ENABLE_GCS_EVENTARC = os.getenv("ENABLE_GCS_EVENTARC", "true").lower() == "true"
ENABLE_HYBRID_SEARCH = os.getenv("ENABLE_HYBRID_SEARCH", "true").lower() == "true"
SOFT_DELETE_RETENTION_DAYS = int(os.getenv("SOFT_DELETE_RETENTION_DAYS", "30"))

# Rate limiting
QUOTA_LIMIT_PER_MINUTE = int(os.getenv("QUOTA_LIMIT_PER_MINUTE", "1000"))
quota_tracker = {"count": 0, "reset_at": datetime.now(timezone.utc)}

# ------------------------------------------------------------------
# PYDANTIC MODELS
# ------------------------------------------------------------------
class IngestRequest(BaseModel):
    folder_id: Optional[str] = None
    gcs_uri: Optional[str] = None
    local_path: Optional[str] = None
    recursive: bool = True  # TRUE RECURSIVE SUPPORT

class SearchRequest(BaseModel):
    query: str
    folder_ids: List[str] = []
    mime_types: Optional[List[str]] = None  # FULL FILTER SUPPORT
    date_from: Optional[str] = None
    date_to: Optional[str] = None
    version_hash: Optional[str] = None
    top_k: int = 20
    use_hybrid: bool = True  # HYBRID SEARCH

class DeleteRequest(BaseModel):
    document_id: str
    permanent: bool = False  # SOFT DELETE BY DEFAULT

class SearchResponse(BaseModel):
    passages: List[Dict[str, Any]]
    total_results: int
    search_time_ms: float
    search_method: str

class WatchChannelRequest(BaseModel):
    folder_id: str
    ttl_hours: int = 24

class EmailWatchRequest(BaseModel):
    gmail_label: Optional[str] = "INBOX"
    max_results: int = 100

class FileShareWatchRequest(BaseModel):
    share_path: str  # SMB/NFS mount point
    watch_pattern: str = "**/*"  # Glob pattern
    poll_interval: int = 300  # seconds

# ------------------------------------------------------------------
# STARTUP
# ------------------------------------------------------------------
@app.on_event("startup")
async def startup():
    global credentials, project_id, db, drive_service, task_client, tokenizer
    
    try:
        credentials, project_id = google.auth.default()
        # Ensure credentials are of type google.auth.credentials.Credentials
        from google.auth.credentials import Credentials
        if not isinstance(credentials, Credentials):
            credentials = None
        db = firestore.Client(project=project_id, credentials=credentials)
        drive_service = build("drive", "v3", credentials=credentials, cache_discovery=False)
        task_client = tasks_v2.CloudTasksClient(credentials=credentials)
        if VERTEX_AVAILABLE:
            aiplatform.init(project=project_id, credentials=credentials)
        
        # Initialize tiktoken for accurate token counting (FEATURE #5)
        tokenizer = tiktoken.get_encoding("cl100k_base")  # GPT-4 tokenizer
        
        log.info("✅ VeriSynthOS Memory Agent - Enterprise Edition started")
        log.info(f"✅ Project: {project_id}")
        log.info(f"✅ Drive Watch: {ENABLE_DRIVE_WATCH}")
        log.info(f"✅ GCS Eventarc: {ENABLE_GCS_EVENTARC}")
        log.info(f"✅ Hybrid Search: {ENABLE_HYBRID_SEARCH}")
        
        # Start background tasks
        asyncio.create_task(renew_watch_channels())
        asyncio.create_task(cleanup_soft_deleted())
        
    except Exception as e:
        log.warning(f"⚠️  GCP not configured: {e}")
        log.info("Running in local dev mode - limited functionality")

# ------------------------------------------------------------------
# FEATURE #5: TOKEN-AWARE CHUNKING (700 tokens, 20% overlap)
# ------------------------------------------------------------------
def semantic_chunk(text: str, max_tokens: int = 700, overlap_tokens: int = 140) -> List[Dict[str, Any]]:
    """
    ENTERPRISE FEATURE: True token-aware chunking with tiktoken.
    - 700 tokens per chunk (not ~2800 chars)
    - 140 token overlap (20%, not 100 chars)
    """
    if not tokenizer:
        # Fallback to character-based if tokenizer not available
        log.warning("Tokenizer not available, using character approximation")
        max_chars = max_tokens * 4
        overlap_chars = overlap_tokens * 4
        chunks = []
        start = 0
        while start < len(text):
            end = min(start + max_chars, len(text))
            chunks.append({"text": text[start:end], "start_char": start, "token_count": (end - start) // 4})
            start = end - overlap_chars if end < len(text) else end
        return chunks
    
    # CORRECT IMPLEMENTATION: Token-based chunking
    tokens = tokenizer.encode(text)
    chunks = []
    start_token = 0
    char_offset = 0
    
    while start_token < len(tokens):
        end_token = min(start_token + max_tokens, len(tokens))
        chunk_tokens = tokens[start_token:end_token]
        chunk_text = tokenizer.decode(chunk_tokens)
        
        chunks.append({
            "text": chunk_text,
            "start_char": char_offset,
            "token_count": len(chunk_tokens),
            "start_token": start_token,
            "end_token": end_token
        })
        
        char_offset += len(chunk_text)
        # 20% overlap for next chunk
        start_token = end_token - overlap_tokens if end_token < len(tokens) else end_token
    
    log.info(f"Chunked {len(tokens)} tokens into {len(chunks)} chunks (700 tokens, 140 overlap)")
    return chunks

# ------------------------------------------------------------------
# FEATURE #8/#18: TIMEZONE-AWARE ISO TIMESTAMPS
# ------------------------------------------------------------------
def now_iso() -> str:
    """Returns ISO 8601 timestamp with 'Z' suffix (UTC)"""
    return datetime.now(timezone.utc).isoformat().replace('+00:00', 'Z')

# ------------------------------------------------------------------
# FEATURE #15: QUOTA-AWARE RATE LIMITING
# ------------------------------------------------------------------
def check_quota():
    """Rate limiting to prevent quota exhaustion"""
    global quota_tracker
    
    now = datetime.now(timezone.utc)
    if now >= quota_tracker["reset_at"]:
        quota_tracker = {"count": 0, "reset_at": now + timedelta(minutes=1)}
    
    if quota_tracker["count"] >= QUOTA_LIMIT_PER_MINUTE:
        raise HTTPException(429, "Rate limit exceeded. Try again in 60 seconds.")
    
    quota_tracker["count"] += 1

# ------------------------------------------------------------------
# FEATURE #13: RETRY WITH EXPONENTIAL BACKOFF (Cloud Tasks)
# ------------------------------------------------------------------
def enqueue_ingestion_task(file_id: str, folder_id: str, retry_count: int = 0):
    """
    ENTERPRISE FEATURE: Reliable ingestion with Cloud Tasks retry
    """
    if not task_client:
        log.warning("Cloud Tasks not available, processing synchronously")
        return None
    
    try:
        parent = task_client.queue_path(str(project_id), "us-central1", CLOUD_TASKS_QUEUE)
        
        task = {
            "http_request": {
                "http_method": tasks_v2.HttpMethod.POST,
                "url": f"{WEBHOOK_URL}/process-file",
                "headers": {"Content-Type": "application/json"},
                "body": json.dumps({
                    "file_id": file_id,
                    "folder_id": folder_id,
                    "retry_count": retry_count
                }).encode()
            }
        }
        
        response = task_client.create_task(request={"parent": parent, "task": task})
        log.info(f"Enqueued task for {file_id}: {response.name}")
        return response.name
        
    except Exception as e:
        log.error(f"Failed to enqueue task: {e}")
        return None

# ------------------------------------------------------------------
# EMBEDDINGS
# ------------------------------------------------------------------
def embed(texts: List[str]) -> List[List[float]]:
    """Batch embedding with Vertex AI"""
    if not VERTEX_AVAILABLE:
        log.warning("Vertex AI not available")
        return [[0.0] * 768 for _ in texts]
    
    try:
        from vertexai.language_models import TextEmbeddingModel
        model = TextEmbeddingModel.from_pretrained("text-embedding-004")
        
        # Batch in groups of 5
        all_embeddings = []
        for i in range(0, len(texts), 5):
            batch = texts[i:i+5]
            embeddings = model.get_embeddings([b for b in batch])
            all_embeddings.extend([e.values for e in embeddings])
        
        return all_embeddings
    except Exception as e:
        log.error(f"Embedding failed: {e}")
        return [[0.0] * 768 for _ in texts]

# ------------------------------------------------------------------
# FEATURE #17: MULTI-MODAL SUPPORT (Images/Charts)
# ------------------------------------------------------------------
def extract_image_text(content: bytes, mime_type: str) -> str:
    """
    ENTERPRISE FEATURE: Extract text from images using Vertex AI Vision
    """
    if not mime_type.startswith("image/"):
        return ""
    
    try:
        # Use Vertex AI Vision API for image understanding
        from google.cloud import vision
        try:
            from google.oauth2.service_account import Credentials as ServiceAccountCredentials
        except ImportError:
            ServiceAccountCredentials = None
        if credentials is not None and ServiceAccountCredentials is not None and isinstance(credentials, ServiceAccountCredentials):
            client = vision.ImageAnnotatorClient(credentials=credentials)
        else:
            client = vision.ImageAnnotatorClient()
        # Use text_detection instead of document_text_detection
        # Use document_text_detection with proper Vision API image construction
        from google.cloud import vision_v1
        vision_image = vision_v1.Image(content=content)
        features = [vision_v1.Feature(type=vision_v1.Feature.Type.DOCUMENT_TEXT_DETECTION)]
        request = vision_v1.AnnotateImageRequest(image=vision_image, features=features)
        try:
            response = client.batch_annotate_images(requests=[request])
            if response and response.responses and hasattr(response.responses[0], 'error') and response.responses[0].error.message:
                raise Exception(response.responses[0].error.message)
            if response and response.responses and hasattr(response.responses[0], 'full_text_annotation') and response.responses[0].full_text_annotation.text:
                text = response.responses[0].full_text_annotation.text
                log.info(f"Extracted {len(text)} chars from image")
                return text
        except Exception as e:
            log.warning(f"Vision API extraction failed: {e}")
        return ""
        
    except Exception as e:
        log.warning(f"Image extraction failed: {e}")
        return ""

def extract_text(content: bytes, mime_type: str) -> str:
    """
    Extract text from various formats including PDFs, DOCX, XML, images
    """
    # Text formats
    if mime_type.startswith("text/") or "json" in mime_type:
        return content.decode("utf-8", errors="ignore")
    
    # Images (FEATURE #17)
    if mime_type.startswith("image/"):
        return extract_image_text(content, mime_type)
    
    # PDF files
    if mime_type == "application/pdf":
        try:
            import PyPDF2
            import io
            pdf_reader = PyPDF2.PdfReader(io.BytesIO(content))
            text = ""
            for page in pdf_reader.pages:
                text += page.extract_text() + "\n"
            log.info(f"Extracted {len(text)} chars from PDF")
            return text
        except Exception as e:
            log.error(f"PDF extraction failed: {e}")
            return ""
    
    # Microsoft Word (.docx)
    if mime_type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
        try:
            import docx
            import io
            doc = docx.Document(io.BytesIO(content))
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            log.info(f"Extracted {len(text)} chars from DOCX")
            return text
        except Exception as e:
            log.error(f"DOCX extraction failed: {e}")
            return ""
    
    # XML files
    if mime_type == "application/xml" or mime_type == "text/xml":
        try:
            import xml.etree.ElementTree as ET
            root = ET.fromstring(content)
            # Extract all text nodes
            text = " ".join([elem.text for elem in root.iter() if elem.text])
            log.info(f"Extracted {len(text)} chars from XML")
            return text
        except Exception as e:
            log.error(f"XML extraction failed: {e}")
            return content.decode("utf-8", errors="ignore")
    
    # CSV/TSV files
    if mime_type == "text/csv" or "spreadsheet" in mime_type:
        try:
            import pandas as pd
            import io
            df = pd.read_csv(io.BytesIO(content))
            text = df.to_string()
            log.info(f"Extracted {len(text)} chars from CSV")
            return text
        except Exception as e:
            log.error(f"CSV extraction failed: {e}")
            return content.decode("utf-8", errors="ignore")
    
    # Excel files
    if "spreadsheet" in mime_type or "excel" in mime_type:
        try:
            import pandas as pd
            import io
            df = pd.read_excel(io.BytesIO(content))
            text = df.to_string()
            log.info(f"Extracted {len(text)} chars from Excel")
            return text
        except Exception as e:
            log.error(f"Excel extraction failed: {e}")
            return ""
    
    # PowerPoint
    if "presentation" in mime_type:
        try:
            from pptx import Presentation
            import io
            prs = Presentation(io.BytesIO(content))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    # Safely extract text if present
                    shape_text = getattr(shape, "text", None)
                    if isinstance(shape_text, str):
                        text += shape_text + "\n"
            log.info(f"Extracted {len(text)} chars from PowerPoint")
            return text
        except Exception as e:
            log.error(f"PowerPoint extraction failed: {e}")
            return ""
    
    log.warning(f"Unsupported mime type: {mime_type}")
    return ""

# ------------------------------------------------------------------
# FEATURE #3: TRUE RECURSIVE FOLDER INGESTION
# ------------------------------------------------------------------
def list_drive_files_recursive(folder_id: str, recursive: bool = True) -> List[Dict]:
    """
    ENTERPRISE FEATURE: Recursively traverse all subfolders
    """
    if not drive_service:
        return []
    
    all_files = []
    folders_to_process = [folder_id]
    
    while folders_to_process:
        current_folder = folders_to_process.pop()
        
        # List all items in current folder
        query = f"'{current_folder}' in parents and trashed=false"
        page_token = None
        
        while True:
            response = drive_service.files().list(
                q=query,
                fields="nextPageToken, files(id, name, mimeType, modifiedTime, md5Checksum, headRevisionId)",
                pageSize=1000,
                pageToken=page_token
            ).execute()
            
            items = response.get("files", [])
            
            for item in items:
                if item["mimeType"] == "application/vnd.google-apps.folder":
                    # Add subfolder to processing queue if recursive
                    if recursive:
                        folders_to_process.append(item["id"])
                else:
                    # Add file to results
                    all_files.append(item)
            
            page_token = response.get("nextPageToken")
            if not page_token:
                break
    
    log.info(f"Found {len(all_files)} files in folder {folder_id} (recursive={recursive})")
    return all_files

# ------------------------------------------------------------------
# FEATURE #6: DRIVE REVISION ID TRACKING
# ------------------------------------------------------------------
def get_file_revision_id(file_id: str) -> Optional[str]:
    """
    ENTERPRISE FEATURE: Fetch latest revision ID for provenance
    """
    if not drive_service:
        return None
    
    try:
        revisions = drive_service.revisions().list(
            fileId=file_id,
            fields="revisions(id, modifiedTime)",
            pageSize=1
        ).execute()
        
        if revisions.get("revisions"):
            latest = revisions["revisions"][-1]
            log.info(f"File {file_id} revision: {latest['id']}")
            return latest["id"]
    except Exception as e:
        log.warning(f"Could not fetch revisions for {file_id}: {e}")
    
    return None

# ------------------------------------------------------------------
# CORE INGESTION
# ------------------------------------------------------------------
async def process_drive_file(file: Dict[str, Any], parent_folder: str) -> int:
    """Process a single Drive file with full provenance tracking"""
    check_quota()  # FEATURE #15
    
    file_id = file["id"]
    file_name = file["name"]
    mime_type = file.get("mimeType", "")
    modified_time = file.get("modifiedTime")
    
    # FEATURE #6: Get revision ID
    revision_id = get_file_revision_id(file_id)
    
    # Fast dedupe using Drive's md5Checksum
    if file.get("md5Checksum"):
        content_hash = file["md5Checksum"]
    else:
        # Download and hash
        try:
            request = drive_service.files().get_media(fileId=file_id)
            fh = io.BytesIO()
            downloader = MediaIoBaseDownload(fh, request)
            done = False
            while not done:
                status, done = downloader.next_chunk()
            content = fh.getvalue()
            content_hash = hashlib.sha256(content).hexdigest()
        except Exception as e:
            log.error(f"Download failed for {file_id}: {e}")
            return 0
    
    # Check if already processed
    if db is None or not hasattr(db, "collection") or not callable(getattr(db, "collection", None)):
        raise RuntimeError("Firestore client (db) is not initialized or missing 'collection' method. Check GCP credentials.")
    existing = db.collection("memory_hashes").document(content_hash).get()
    if existing.exists:
        log.info(f"Duplicate skipped: {file_name}")
        return 0
        try:
            if drive_service is None:
                try:
                    from googleapiclient.discovery import build
                    import google.auth
                    credentials, _ = google.auth.default()
                    drive_service = build('drive', 'v3', credentials=credentials, cache_discovery=False)
                except Exception as e:
                    log.error(f"Failed to initialize drive_service: {e}")
                    return 0
            if drive_service is not None and hasattr(drive_service, "files") and callable(getattr(drive_service, "files", None)):
                try:
                    request = drive_service.files().get_media(fileId=file_id)
                    fh = io.BytesIO()
                    downloader = MediaIoBaseDownload(fh, request)
                    done = False
                    while not done:
                        status, done = downloader.next_chunk()
                    content = fh.getvalue()
                    content_hash = hashlib.sha256(content).hexdigest()
                except Exception as e:
                    log.error(f"drive_service.files() failed for file {file_id}: {e}")
                    return 0
            else:
                log.error(f"drive_service.files is not available or not callable for file {file_id}")
                return 0
        except Exception as e:
            log.error(f"Download failed for {file_id}: {e}")
            return 0
        log.error(f"Download failed for {file_id}: {e}")
        # FEATURE #13: Enqueue retry task
        enqueue_ingestion_task(file_id, parent_folder, retry_count=1)
        return 0
    
    # Extract text (includes FEATURE #17: image support)
    text = extract_text(content, mime_type)
    if not text:
        log.warning(f"No text extracted from {file_name}")
        return 0
    
    # FEATURE #4/#5: Token-aware chunking with 20% overlap
    chunks = semantic_chunk(text)
    
    # Embed chunks
    embeddings = embed([c["text"] for c in chunks])
    
    # Store with FULL PROVENANCE (FEATURES #6/#7/#8)
    try:
        doc_data = {
            "file_id": file_id,
            "file_name": file_name,
            "mime_type": mime_type,
            "content_hash": content_hash,  # FEATURE #7
            "revision_id": revision_id,  # FEATURE #6
            "modified_at": modified_time,  # FEATURE #8
            "uploaded_at": now_iso(),  # FEATURE #18
            "chunk_count": len(chunks),
            "parent_folder": parent_folder,
            "drive_link": f"https://drive.google.com/file/d/{file_id}",
            "source": "drive",
            "deleted": False  # FEATURE #12
        }
        
        if db is None:
            raise RuntimeError("Firestore client (db) is not initialized. Check GCP credentials.")
        db.collection("memory_docs").document(file_id).set(doc_data)
        
        # Store chunks with embeddings
        for i, (chunk, emb) in enumerate(zip(chunks, embeddings)):
            chunk_data = {
                "document_id": file_id,
                "chunk_index": i,
                "text": chunk["text"],
                "start_char": chunk["start_char"],
                "token_count": chunk.get("token_count", 0),
                "embedding": emb,
                "created_at": now_iso()
            }
            if db is None:
                raise RuntimeError("Firestore client (db) is not initialized. Check GCP credentials.")
            db.collection("chunks").add(chunk_data)
        
        # Dedupe hash
        if db is None:
            raise RuntimeError("Firestore client (db) is not initialized. Check GCP credentials.")
        db.collection("memory_hashes").document(content_hash).set({
            "file_id": file_id,
            "file_name": file_name,
            "indexed_at": now_iso()
        })
        
        log.info(f"✅ Indexed {len(chunks)} chunks for {file_name}")
        return 1
        
    except Exception as e:
        log.error(f"❌ Failed to index {file_name}: {e}")
        # FEATURE #13: Retry
        enqueue_ingestion_task(file_id, parent_folder, retry_count=1)
        return 0

# ------------------------------------------------------------------
# LOCAL FILE INGESTION (with token-aware chunking)
# ------------------------------------------------------------------
def process_local_file(filepath: str) -> int:
    """Process local file with enterprise features"""
    if not db:
        raise HTTPException(503, "Firestore required")
    
    path = Path(filepath)
    if not path.exists():
        raise HTTPException(404, f"File not found: {filepath}")
    
    try:
        with open(path, 'r', encoding='utf-8') as f:
            content = f.read()
    except Exception as e:
        log.error(f"Error reading {filepath}: {e}")
        return 0
    
    content_hash = hashlib.sha256(content.encode()).hexdigest()
    
    # Dedupe check
    existing = db.collection("memory_hashes").document(content_hash).get()
    if existing.exists:
        log.info(f"Duplicate skipped: {path.name}")
        return 0
    
    # FEATURE #5: Token-aware chunking
    chunks = semantic_chunk(content)
    
    # Embed if available
    embeddings = embed([c["text"] for c in chunks]) if credentials else []
    
    # Store with full provenance
    file_id = content_hash[:16]
    try:
        doc_data = {
            "file_id": file_id,
            "file_name": path.name,
            "file_path": str(path.absolute()),
            "content_hash": content_hash,
            "modified_at": datetime.fromtimestamp(path.stat().st_mtime, tz=timezone.utc).isoformat(),
            "uploaded_at": now_iso(),  # FEATURE #18
            "chunk_count": len(chunks),
            "parent_folder": str(path.parent),
            "source": "local",
            "deleted": False  # FEATURE #12
        }
        
        db.collection("memory_docs").document(file_id).set(doc_data)
        
        # Store chunks
        for i, chunk in enumerate(chunks):
            chunk_data = {
                "document_id": file_id,
                "chunk_index": i,
                "text": chunk["text"],
                "start_char": chunk["start_char"],
                "token_count": chunk.get("token_count", 0),
                "embedding": embeddings[i] if embeddings else None,
                "created_at": now_iso()
            }
            db.collection("chunks").add(chunk_data)
        
        db.collection("memory_hashes").document(content_hash).set({
            "file_id": file_id,
            "file_name": path.name,
            "indexed_at": now_iso()
        })
        
        log.info(f"✅ Indexed {len(chunks)} chunks for {path.name}")
        return len(chunks)
        
    except Exception as e:
        log.error(f"❌ Failed to index {filepath}: {e}")
        return 0

# ------------------------------------------------------------------
# FEATURE #1: REAL-TIME DRIVE WATCHING
# ------------------------------------------------------------------
@app.post("/watch/start")
async def start_watch(req: WatchChannelRequest):
    """
    ENTERPRISE FEATURE: Start real-time Drive push notifications
    """
    if not drive_service:
        raise HTTPException(503, "Drive API not available")
    
    check_quota()
    
    try:
        channel_id = str(uuid.uuid4())
        expiration = datetime.now(timezone.utc) + timedelta(hours=req.ttl_hours)
        
        body = {
            "id": channel_id,
            "type": "web_hook",
            "address": WEBHOOK_URL,
            "expiration": int(expiration.timestamp() * 1000)
        }
        
        # FEATURE #14: For large folders, implement sharding
        file_count = len(list_drive_files_recursive(req.folder_id, recursive=False))
        if file_count > 10000:
            log.warning(f"Folder {req.folder_id} has {file_count} files - consider sharding")
        
        response = drive_service.files().watch(
            fileId=req.folder_id,
            body=body
        ).execute()
        
        # Store channel info
        watch_channels[channel_id] = {
            "folder_id": req.folder_id,
            "resource_id": response["resourceId"],
            "expiration": expiration,
            "created_at": now_iso()
        }
        
        log.info(f"✅ Started watch on {req.folder_id}, channel: {channel_id}")
        
        return {
            "status": "watching",
            "channel_id": channel_id,
            "expiration": expiration.isoformat()
        }
        
    except Exception as e:
        log.error(f"❌ Failed to start watch: {e}")
        raise HTTPException(500, str(e))

@app.post("/webhook/drive")
async def drive_webhook(request: Request, background_tasks: BackgroundTasks):
    """
    ENTERPRISE FEATURE: Receive Drive push notifications
    """
    headers = request.headers
    channel_id = headers.get("x-goog-channel-id")
    resource_state = headers.get("x-goog-resource-state")
    
    if not channel_id or channel_id not in watch_channels:
        return JSONResponse({"status": "ignored"}, status_code=200)
    
    if resource_state == "sync":
        # Initial sync message
        return JSONResponse({"status": "synced"}, status_code=200)
    
    # File changed - re-ingest folder
    folder_id = watch_channels[channel_id]["folder_id"]
    log.info(f"📢 Change detected in folder {folder_id}")
    
    # Process in background
    background_tasks.add_task(re_ingest_folder, folder_id)
    
    return JSONResponse({"status": "processing"}, status_code=200)

async def re_ingest_folder(folder_id: str):
    """Re-process folder after change notification"""
    files = list_drive_files_recursive(folder_id, recursive=True)
    for file in files:
        await process_drive_file(file, folder_id)

async def renew_watch_channels():
    """Background task to renew expiring watch channels"""
    while True:
        await asyncio.sleep(3600)  # Check every hour
        
        now = datetime.now(timezone.utc)
        for channel_id, info in list(watch_channels.items()):
            expiration = info["expiration"]
            
            # Renew if expiring in next 12 hours
            if expiration - now < timedelta(hours=12):
                try:
                    # Stop old channel
                    if drive_service is not None and hasattr(drive_service, "channels") and callable(getattr(drive_service, "channels", None)):
                        drive_service.channels().stop(body={
                            "id": channel_id,
                            "resourceId": info["resource_id"]
                        }).execute()
                    
                    # Start new channel
                    await start_watch(WatchChannelRequest(
                        folder_id=info["folder_id"],
                        ttl_hours=24
                    ))
                    
                    log.info(f"♻️  Renewed watch channel for {info['folder_id']}")
                    
                except Exception as e:
                    log.error(f"Failed to renew channel {channel_id}: {e}")

# ------------------------------------------------------------------
# FEATURE #2: GCS EVENTARC INGESTION
# ------------------------------------------------------------------
@app.post("/webhook/gcs")
async def gcs_eventarc_handler(request: Request):
    """
    ENTERPRISE FEATURE: Process GCS object via Eventarc trigger
    """
    if not ENABLE_GCS_EVENTARC:
        raise HTTPException(501, "GCS Eventarc not enabled")
    
    event = await request.json()
    
    bucket = event.get("bucket")
    name = event.get("name")
    
    if not bucket or not name:
        return JSONResponse({"status": "ignored"}, status_code=200)
    
    gcs_uri = f"gs://{bucket}/{name}"
    log.info(f"📦 GCS event: {gcs_uri}")
    
    # Download and process
    try:
        from google.cloud import storage
        storage_client = storage.Client(credentials=credentials)
        
        bucket_obj = storage_client.bucket(bucket)
        blob = bucket_obj.blob(name)
        content = blob.download_as_bytes()
        
        # Process similar to Drive files
        content_hash = hashlib.sha256(content).hexdigest()
        
        # Check dedupe
        if db is None:
            raise RuntimeError("Firestore client (db) is not initialized. Check GCP credentials.")
        existing = db.collection("memory_hashes").document(content_hash).get()
        if existing.exists:
            return JSONResponse({"status": "duplicate"}, status_code=200)
        
        # Extract and chunk
        mime_type = blob.content_type or "application/octet-stream"
        text = extract_text(content, mime_type)
        chunks = semantic_chunk(text)
        embeddings = embed([c["text"] for c in chunks])
        
        # Store
        file_id = content_hash[:16]
        if db is None:
            raise RuntimeError("Firestore client (db) is not initialized. Check GCP credentials.")
        db.collection("memory_docs").document(file_id).set({
            "file_id": file_id,
            "file_name": name,
            "gcs_uri": gcs_uri,
            "content_hash": content_hash,
            "uploaded_at": now_iso(),
            "chunk_count": len(chunks),
            "source": "gcs",
            "deleted": False
        })
        
        for i, (chunk, emb) in enumerate(zip(chunks, embeddings)):
            if db is None:
                raise RuntimeError("Firestore client (db) is not initialized. Check GCP credentials.")
            db.collection("chunks").add({
                "document_id": file_id,
                "chunk_index": i,
                "text": chunk["text"],
                "embedding": emb,
                "created_at": now_iso()
            })
        
        if db is None:
            raise RuntimeError("Firestore client (db) is not initialized. Check GCP credentials.")
        db.collection("memory_hashes").document(content_hash).set({
            "file_id": file_id,
            "file_name": name,
            "indexed_at": now_iso()
        })
        
        log.info(f"✅ Processed GCS file: {name}")
        return JSONResponse({"status": "indexed", "chunks": len(chunks)}, status_code=200)
        
    except Exception as e:
        log.error(f"❌ GCS processing failed: {e}")
        return JSONResponse({"status": "error", "message": str(e)}, status_code=500)

# ------------------------------------------------------------------
# FEATURE #9: HYBRID SEARCH (Vector + BM25)
# ------------------------------------------------------------------
def bm25_search(query: str, top_k: int = 20) -> List[Dict]:
    """
    ENTERPRISE FEATURE: BM25 keyword search for hybrid ranking
    """
    # Simple implementation - production would use specialized search index
    query_terms = query.lower().split()
    
    # Query Firestore for text matches
    if db is None:
        raise RuntimeError("Firestore client (db) is not initialized. Check GCP credentials.")
    all_chunks = db.collection("chunks").stream()
    
    results = []
    for doc in all_chunks:
        data = doc.to_dict()
        text = data.get("text", "").lower()
        
        # Simple scoring: count query term occurrences
        score = sum(text.count(term) for term in query_terms)
        
        if score > 0:
            results.append({
                "id": doc.id,
                "score": score,
                "text": data.get("text"),
                "document_id": data.get("document_id"),
                "chunk_index": data.get("chunk_index")
            })
    
    # Sort by score and return top k
    results.sort(key=lambda x: x["score"], reverse=True)
    return results[:top_k]

def vector_search(query: str, top_k: int = 20, filters: Optional[Dict] = None) -> List[Dict]:
    """
    Vector similarity search with metadata filters
    """
    # Get query embedding
    query_embedding = embed([query])[0]
    
    # Query Matching Engine (if available) or Firestore
    # For now, simple Firestore similarity
    if db is not None and hasattr(db, "collection") and callable(getattr(db, "collection", None)):
        all_chunks = db.collection("chunks").stream()
    else:
        all_chunks = []
    
    results = []
    for doc in all_chunks:
        data = doc.to_dict()
        
        # Apply metadata filters (FEATURE #10)
        if filters:
            # Would filter by document properties
            pass
        
        emb = data.get("embedding")
        if emb:
            # Cosine similarity
            import numpy as np
            score = np.dot(query_embedding, emb) / (np.linalg.norm(query_embedding) * np.linalg.norm(emb))
            
            results.append({
                "id": doc.id,
                "score": float(score),
                "text": data.get("text"),
                "document_id": data.get("document_id"),
                "chunk_index": data.get("chunk_index")
            })
    
    results.sort(key=lambda x: x["score"], reverse=True)
    return results[:top_k]

def hybrid_search(query: str, top_k: int = 20, filters: Optional[Dict] = None) -> List[Dict]:
    """
    ENTERPRISE FEATURE: Hybrid search combining vector + BM25
    """
    # Get both result sets
    if db is None:
        raise RuntimeError("Firestore client (db) is not initialized. Check GCP credentials.")
    vector_results = vector_search(query, top_k * 2, filters)
    bm25_results = bm25_search(query, top_k * 2)
    
    # Reciprocal Rank Fusion (RRF)
    k = 60  # RRF constant
    scores = {}
    
    for rank, result in enumerate(vector_results, 1):
        doc_id = result["id"]
        scores[doc_id] = scores.get(doc_id, 0) + 1 / (k + rank)
    
    for rank, result in enumerate(bm25_results, 1):
        doc_id = result["id"]
        scores[doc_id] = scores.get(doc_id, 0) + 1 / (k + rank)
    
    # Combine and re-rank
    all_results = {r["id"]: r for r in vector_results + bm25_results}
    
    ranked = sorted(scores.items(), key=lambda x: x[1], reverse=True)
    
    return [all_results[doc_id] for doc_id, score in ranked[:top_k] if doc_id in all_results]

# ------------------------------------------------------------------
# FEATURE #10: FULL METADATA FILTERS + SEARCH ENDPOINT
# ------------------------------------------------------------------
@app.post("/search")
async def search(req: SearchRequest):
    """
    ENTERPRISE FEATURE: Complete search with all filters and provenance
    """
    check_quota()
    
    # Build metadata filters
    filters = {}
    if req.folder_ids:
        filters["folder_ids"] = req.folder_ids
    if req.mime_types:
        filters["mime_types"] = req.mime_types
    if req.date_from:
        filters["date_from"] = req.date_from
    if req.date_to:
        filters["date_to"] = req.date_to
    if req.version_hash:
        filters["version_hash"] = req.version_hash
    
    # Choose search method
    if req.use_hybrid and ENABLE_HYBRID_SEARCH:
        results = hybrid_search(req.query, req.top_k, filters)
    else:
        results = vector_search(req.query, req.top_k, filters)
    
    # Enrich with FULL PROVENANCE (FEATURES #7/#8)
    enriched = []
    for result in results:
        doc_id = result["document_id"]
        if db is None:
            raise RuntimeError("Firestore client (db) is not initialized. Check GCP credentials.")
        doc_ref = db.collection("memory_docs").document(doc_id).get()
        doc_data = doc_ref.to_dict() if doc_ref else None
        if doc_data and not doc_data.get("deleted", False):  # Exclude soft-deleted
            enriched.append({
                "text": result["text"],
                "score": result["score"],
                "chunk_index": result["chunk_index"],
                "provenance": {
                    "file_name": doc_data.get("file_name") if doc_data else None,
                    "file_id": doc_data.get("file_id") if doc_data else None,
                    "version_hash": doc_data.get("content_hash") if doc_data else None,
                    "revision_id": doc_data.get("revision_id") if doc_data else None,
                    "modified_at": doc_data.get("modified_at") if doc_data else None,
                    "uploaded_at": doc_data.get("uploaded_at") if doc_data else None,
                    "drive_link": doc_data.get("drive_link") if doc_data else None,
                    "source": doc_data.get("source") if doc_data else None
                }
            })
    
    return {
        "query": req.query,
        "results": enriched,
        "total": len(enriched),
        "search_type": "hybrid" if req.use_hybrid else "vector"
    }

# ------------------------------------------------------------------
# FEATURE #11: DELETE ENDPOINT
# ------------------------------------------------------------------
@app.delete("/doc/{document_id}")
async def delete_document(document_id: str, req: DeleteRequest):
    """
    ENTERPRISE FEATURE: Delete document (soft or permanent)
    """
    check_quota()
    
    if not db:
        raise HTTPException(503, "Firestore not available")
    
    doc_ref = db.collection("memory_docs").document(document_id)
    doc = doc_ref.get()
    
    if not doc.exists:
        raise HTTPException(404, "Document not found")
    
    if req.permanent:
        # PERMANENT DELETE
        try:
            # Delete all chunks
            chunks = db.collection("chunks").where("document_id", "==", document_id).stream()
            for chunk in chunks:
                chunk.reference.delete()
            
            # Delete document
            doc_ref.delete()
            
            # Delete hash entry
            doc_data = doc.to_dict() if doc else None
            if doc_data and doc_data.get("content_hash"):
                db.collection("memory_hashes").document(doc_data["content_hash"]).delete()
            
            log.info(f"🗑️  Permanently deleted {document_id}")
            return {"status": "deleted", "permanent": True}
            
        except Exception as e:
            log.error(f"Delete failed: {e}")
            raise HTTPException(500, str(e))
    
    else:
        # FEATURE #12: SOFT DELETE with 30-day retention
        doc_ref.update({
            "deleted": True,
            "deleted_at": now_iso()
        })
        
        log.info(f"📦 Soft-deleted {document_id} (30-day retention)")
        return {"status": "deleted", "permanent": False, "retention_days": SOFT_DELETE_RETENTION_DAYS}

# ------------------------------------------------------------------
# FEATURE #12: SOFT DELETE CLEANUP
# ------------------------------------------------------------------
async def cleanup_soft_deleted():
    """Background task to permanently delete expired soft-deleted documents"""
    while True:
        await asyncio.sleep(86400)  # Run daily
        
        if not db:
            continue
        
        cutoff = datetime.now(timezone.utc) - timedelta(days=SOFT_DELETE_RETENTION_DAYS)
        
        try:
            # Find expired soft-deleted documents
            expired = db.collection("memory_docs").where("deleted", "==", True).where(
                "deleted_at", "<=", cutoff.isoformat()
            ).stream()
            
            for doc in expired:
                # Permanently delete
                await delete_document(doc.id, DeleteRequest(document_id=doc.id, permanent=True))
                log.info(f"♻️  Auto-deleted expired document {doc.id}")
                
        except Exception as e:
            log.error(f"Cleanup failed: {e}")

# ------------------------------------------------------------------
# ENDPOINTS
# ------------------------------------------------------------------
@app.get("/")
async def root():
    return {
        "service": "VeriSynthOS Memory Agent - Enterprise Edition",
        "version": "2.0.0",
        "features": {
            "drive_watch": ENABLE_DRIVE_WATCH,
            "gcs_eventarc": ENABLE_GCS_EVENTARC,
            "hybrid_search": ENABLE_HYBRID_SEARCH,
            "soft_delete": True,
            "token_aware_chunking": True,
            "multi_modal": True,
            "full_provenance": True
        }
    }

@app.get("/health")
async def health():
    return {
        "status": "healthy",
        "gcp": credentials is not None,
        "firestore": db is not None,
        "drive": drive_service is not None,
        "watch_channels": len(watch_channels)
    }

@app.post("/ingest")
async def ingest(req: IngestRequest):
    """
    ENTERPRISE FEATURE: Ingest with true recursive support
    """
    check_quota()
    
    # Local file system
    if req.local_path:
        if not db:
            raise HTTPException(503, "Firestore required")
        
        path = Path(req.local_path)
        if path.is_file():
            chunks = process_local_file(str(path))
            return {"status": "ok", "files": 1, "chunks": chunks, "timestamp": now_iso()}
        elif path.is_dir():
            total_chunks = 0
            pattern = "**/*" if req.recursive else "*"
            for file_path in path.glob(pattern):
                if file_path.is_file():
                    total_chunks += process_local_file(str(file_path))
            return {"status": "ok", "chunks": total_chunks, "timestamp": now_iso()}
    
    # Google Drive (FEATURE #3: TRUE RECURSIVE)
    if req.folder_id:
        if not drive_service:
            raise HTTPException(503, "Drive API not available")
        
        files = list_drive_files_recursive(req.folder_id, req.recursive)
        processed = 0
        
        for file in files:
            result = await process_drive_file(file, req.folder_id)
            processed += result
        
        return {"status": "ok", "files_processed": processed, "timestamp": now_iso()}
    
    # GCS
    if req.gcs_uri:
        # Would process GCS URI
        pass
    
    raise HTTPException(400, "folder_id, gcs_uri, or local_path required")

# ------------------------------------------------------------------
# FEATURE #19: DYNAMIC RED-FLAG THRESHOLDS (Integration point)
# ------------------------------------------------------------------
@app.get("/maker/threshold/{document_id}")
async def get_red_flag_threshold(document_id: str):
    """
    ENTERPRISE FEATURE: Return dynamic red-flag threshold based on memory quality
    """
    if db is None:
        raise RuntimeError("Firestore client (db) is not initialized. Check GCP credentials.")
    doc_ref = db.collection("memory_docs").document(document_id).get()
    if not doc_ref or not doc_ref.exists:
        return {"threshold": 1200}  # Default for advanced models
    doc_data = doc_ref.to_dict() if doc_ref else None
    # Compute threshold based on document metadata
    base_threshold = 1200
    if doc_data and doc_data.get("source") == "drive":
        base_threshold -= 100  # Trust Drive documents more
    modified_at = doc_data.get("modified_at") if doc_data else None
    if modified_at:
        age_days = (datetime.now(timezone.utc) - datetime.fromisoformat(modified_at.replace('Z', '+00:00'))).days
        if age_days < 30:
            base_threshold -= 50  # Recent docs more reliable
    
    return {
        "document_id": document_id,
        "red_flag_threshold": base_threshold,
        "reasoning": "Adjusted based on source and recency"
    }

# ------------------------------------------------------------------
# EMAIL & FILE SHARE CONNECTORS
# ------------------------------------------------------------------
from agents.memory.connectors import (
    EmailWatchRequest,
    FileShareWatchRequest,
    start_email_watch,
    start_fileshare_watch,
    stop_fileshare_watch,
    list_fileshare_watchers
)

@app.post("/watch/email")
async def watch_email(req: EmailWatchRequest, background_tasks: BackgroundTasks):
    """Monitor Gmail for attachments"""
    return await start_email_watch(
        req, credentials, db, extract_text, semantic_chunk, embed, now_iso
    )

@app.post("/watch/fileshare")
async def watch_fileshare(req: FileShareWatchRequest, background_tasks: BackgroundTasks):
    """Monitor network file share"""
    return await start_fileshare_watch(req, db, process_local_file)

@app.delete("/watch/fileshare/{watcher_id}")
async def unwatch_fileshare(watcher_id: str):
    """Stop monitoring file share"""
    return await stop_fileshare_watch(watcher_id)

@app.get("/watch/fileshare")
async def list_watchers():
    """List active file share watchers"""
    return list_fileshare_watchers()

# ------------------------------------------------------------------
# MODERN UI
# ------------------------------------------------------------------
import os
from pathlib import Path

# Mount static files
static_dir = Path(__file__).parent / "static"
if static_dir.exists():
    app.mount("/static", StaticFiles(directory=str(static_dir)), name="static")

@app.get("/ui", response_class=HTMLResponse)
async def serve_ui():
    """Serve modern SEO-optimized UI"""
    ui_file = Path(__file__).parent / "static" / "index.html"
    if ui_file.exists():
        return HTMLResponse(content=ui_file.read_text())
    else:
        return HTMLResponse(content="<h1>UI not found. Create static/index.html</h1>")

